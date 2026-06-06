"""Generates Verispect-Pitch-Deck.pptx — 15-slide investor/sales deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PURPLE = RGBColor(0x7C,0x6E,0xFF); DARK = RGBColor(0x11,0x18,0x27)
GRAY = RGBColor(0x6B,0x72,0x80); LT = RGBColor(0xED,0xE9,0xFE)
WHITE = RGBColor(0xFF,0xFF,0xFF); GREEN = RGBColor(0x10,0xB9,0x81)
BG = RGBColor(0x0F,0x11,0x17)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def box(s, x, y, w, h):
    return s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

def txt(tf, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, italic=False, space=6):
    p = tf.paragraphs[0] if (len(tf.paragraphs)==1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = "Calibri"
    return p

def bullets(tf, items, size=18, color=DARK, space=10):
    first = True
    for it in items:
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
        first = False
        p.space_after = Pt(space)
        r = p.add_run(); r.text = "•  " + it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name="Calibri"

def accent(s, y=1.15):
    b = s.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(1.6), Inches(0.08))
    b.fill.solid(); b.fill.fore_color.rgb = PURPLE; b.line.fill.background(); b.shadow.inherit=False

def title_only(s, title, sub=None):
    t = box(s, 0.8, 0.45, 11.7, 1.0); txt(t.text_frame, title, 30, DARK, bold=True)
    accent(s)
    if sub:
        u = box(s, 0.8, 1.25, 11.7, 0.6); txt(u.text_frame, sub, 15, GRAY, italic=True)

# 1 — COVER
s = slide(BG)
t = box(s, 0.9, 2.5, 11.5, 2.2)
txt(t.text_frame, "Verispect", 60, WHITE, bold=True)
txt(t.text_frame, "Active AI bias & drift detection + EU AI Act compliance evidence — in one line of code.", 20, RGBColor(0xC8,0xC3,0xF0))
b = s.shapes.add_shape(1, Inches(0.95), Inches(4.15), Inches(2.0), Inches(0.08))
b.fill.solid(); b.fill.fore_color.rgb = PURPLE; b.line.fill.background(); b.shadow.inherit=False
f = box(s, 0.9, 6.4, 11.5, 0.6)
txt(f.text_frame, "Verify + Inspect. Every model. Every call. Every month.   ·   verispectai.com", 13, GRAY)

# 2 — PROBLEM
s = slide(); title_only(s, "The most dangerous AI bug is the one where nothing changed")
t = box(s, 0.8, 2.0, 11.7, 4.5)
bullets(t.text_frame, [
    "Companies put LLMs into decisions that change lives — hiring, credit, triage, legal review.",
    "Providers silently update models behind a stable API name. Behaviour shifts overnight, no changelog.",
    "Fine-tunes degrade. Prompts drift. A model that was fair last month quietly starts discriminating.",
    "Logging tools can't catch it — the inputs and outputs still look normal.",
    "Teams find out when a regulator, an enterprise buyer, or a journalist asks: \"prove your AI behaves.\"",
    "They have no evidence. Building it in-house is a multi-month ML + compliance project nobody has time for.",
], size=18, space=14)

# 3 — INSIGHT
s = slide(BG)
t = box(s, 0.9, 2.6, 11.5, 2.0)
txt(t.text_frame, "You can't detect drift by watching.", 34, WHITE, bold=True)
txt(t.text_frame, "You have to interrogate.", 34, PURPLE, bold=True)
u = box(s, 0.9, 4.7, 11.5, 1.6)
txt(u.text_frame, "Fire calibrated probes — identical candidate profiles differing by one protected characteristic — "
    "at the live model. Measure the gap vs a recorded baseline. Divergence becomes a number you can put in a report.", 17, RGBColor(0xC8,0xC3,0xF0))

# 4 — SOLUTION
s = slide(); title_only(s, "Verispect: Active AI Behavioural Assurance")
cards = [("Active, not passive","Calibrated probes + replays of your own traffic test behaviour continuously — catching change the moment it happens."),
         ("Audit-ready by default","Branded report mapped article-by-article to EU AI Act 9/10/13/14/72 + Annex III. The document buyers and auditors accept."),
         ("Privacy by architecture","Only SHA-256 hashes + embedding vectors leave the customer. We reduce their data-protection surface.")]
x = 0.8
for title, bodytx in cards:
    c = s.shapes.add_shape(1, Inches(x), Inches(2.2), Inches(3.85), Inches(3.6))
    c.fill.solid(); c.fill.fore_color.rgb = LT; c.line.color.rgb = PURPLE; c.line.width=Pt(1); c.shadow.inherit=False
    tf = c.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.25); tf.margin_right=Inches(0.25); tf.margin_top=Inches(0.3)
    txt(tf, title, 19, PURPLE, bold=True, space=10)
    txt(tf, bodytx, 14, DARK)
    x += 4.07

# 5 — HOW IT WORKS
s = slide(); title_only(s, "How it works", "One line. Zero added latency. Three steps.")
steps=[("1","Add one line","client = wrap(OpenAI(...), verispect_key=\"vs_live_...\")  — or swap the base_url."),
       ("2","We probe & score","On a sample of live traffic, Verispect fires bias/consistency probes and scores drift vs your baseline — deterministically, no LLM in the loop."),
       ("3","You get the report","A continuously-updated, audit-ready PDF mapped to the EU AI Act. Send it to your buyer or auditor.")]
y=2.3
for num,h,b in steps:
    c=s.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(0.8), Inches(0.8))
    c.fill.solid(); c.fill.fore_color.rgb=PURPLE; c.line.fill.background(); c.shadow.inherit=False
    ctf=c.text_frame; ctf.word_wrap=True; p=ctf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=num; r.font.size=Pt(26); r.font.bold=True; r.font.color.rgb=WHITE
    t=box(s,1.9,y-0.05,10.6,1.1); txt(t.text_frame,h,19,DARK,bold=True,space=2); txt(t.text_frame,b,14,GRAY)
    y+=1.45

# 6 — WHY NOW
s = slide(); title_only(s, "Why now", "Regulation turns a nice-to-have into a deadline")
t=box(s,0.8,2.0,11.7,4.3)
bullets(t.text_frame,[
    "EU AI Act high-risk (Annex III) obligations operative 2 Aug 2026 — employment/recruitment AI is explicitly high-risk.",
    "Digital Omnibus may move standalone Annex III to 2 Dec 2027 (not yet law) — the obligation still arrives.",
    "Continuous post-market monitoring (Art. 72), logging (Art. 13), bias governance (Art. 10) become mandatory.",
    "Enterprise procurement already asks AI-governance questions today — demand precedes the law.",
    "Parallel pull from Singapore (IMDA AI Verify) and UAE governance frameworks.",
], size=18, space=14)

# 7 — MARKET
s = slide(); title_only(s, "Market", "Positioned on the fastest-growing demand vector: compliance-driven assurance")
stats=[("$2–3.2B","LLM observability market, 2025"),("25–36%","CAGR toward $9–24B by 2030–34"),("~$2.7B","Compliance-assurance TAM (2030, est.)"),("~€1M ARR","3-yr SOM from ~150 right accounts")]
x=0.8
for big,lab in stats:
    c=s.shapes.add_shape(1,Inches(x),Inches(2.4),Inches(2.85),Inches(2.4))
    c.fill.solid(); c.fill.fore_color.rgb=LT; c.line.fill.background(); c.shadow.inherit=False
    tf=c.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    txt(tf,big,30,PURPLE,bold=True,align=PP_ALIGN.CENTER,space=6)
    txt(tf,lab,13,DARK,align=PP_ALIGN.CENTER)
    x+=3.05
u=box(s,0.8,5.2,11.7,1.0)
txt(u.text_frame,"Analysts cite regulatory scrutiny (EU AI Act, NIST AI RMF) as the #1 growth driver — exactly our slice.",14,GRAY,italic=True)

# 8 — COMPETITION
s = slide(); title_only(s, "Competitive landscape", "Observability watches. Assurance verifies.")
rows=[["","Logs","Active probing","Bias probes","AI Act report","No data access"],
      ["Helicone","✓","—","—","—","—"],
      ["LangSmith","✓","—","—","—","—"],
      ["Braintrust","✓","—","—","—","—"],
      ["Verispect","✓","✓","✓ (8)","✓","✓"]]
tbl = s.shapes.add_table(len(rows),len(rows[0]),Inches(0.8),Inches(2.2),Inches(11.7),Inches(3.4)).table
for ci,_ in enumerate(rows[0]):
    tbl.columns[ci].width=Inches(11.7/6)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=val
        pr=cell.text_frame.paragraphs[0]; pr.alignment=PP_ALIGN.CENTER
        rr=pr.runs[0] if pr.runs else pr.add_run()
        rr.font.size=Pt(14); rr.font.name="Calibri"
        if ri==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=DARK; rr.font.color.rgb=WHITE; rr.font.bold=True
        elif row[0]=="Verispect":
            cell.fill.solid(); cell.fill.fore_color.rgb=PURPLE; rr.font.color.rgb=WHITE; rr.font.bold=True
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb=WHITE; rr.font.color.rgb=DARK

# 9 — MOAT
s = slide(); title_only(s, "Why we win — and keep winning")
t=box(s,0.8,2.0,11.7,4.3)
bullets(t.text_frame,[
    "Probe library as IP — calibrated, regulation-mapped probes across 8 protected characteristics. A weekend to copy the proxy; a year to copy the library.",
    "Compliance output as lock-in — the audit history accumulates inside Verispect. Switching means losing evidence right before an audit.",
    "Privacy architecture as trust wedge — \"we literally can't see your data\" closes security review faster than any competitor.",
    "Regulatory tailwind — the buying trigger is a legal deadline, not a preference.",
], size=18, space=16)

# 10 — PRODUCT (status)
s = slide(); title_only(s, "Product — already built", "Working full stack, not a concept")
t=box(s,0.8,2.0,11.7,4.3)
bullets(t.text_frame,[
    "One-line SDK + middleware proxy · zero added latency (async background work).",
    "Privacy-preserving ingest — hashes + embeddings only; golden probes stored on the customer's machine.",
    "20-probe regulatory library + golden-probe replay of real traffic.",
    "Deterministic drift scoring (all-MiniLM-L6-v2 cosine similarity).",
    "Branded EU AI Act compliance PDF — article-by-article mapping, drift log, remediation.",
    "React dashboard, JWT auth, API-key management, multi-tenant.",
], size=17, space=12)

# 11 — BUSINESS MODEL
s = slide(); title_only(s, "Business model", "One product, zero config. Premium, value-based. ~95% gross margin.")
tiers=[("Free","$0","AI Act snapshot — lead magnet"),("Verispect","$1,500/mo","zero config · everything · always-current audit pack"),
       ("Founding 20","$1,500 locked","for life, before it rises to $2,500"),("Enterprise","from $2,500/mo","multi-entity · SSO · SLA · audit support")]
x=0.8
for name,price,desc in tiers:
    c=s.shapes.add_shape(1,Inches(x),Inches(2.4),Inches(2.85),Inches(3.0))
    c.fill.solid(); c.fill.fore_color.rgb=(PURPLE if name=="Verispect" else LT); c.line.fill.background(); c.shadow.inherit=False
    tf=c.text_frame; tf.word_wrap=True; tf.margin_top=Inches(0.3)
    col = WHITE if name=="Verispect" else DARK
    txt(tf,name,20,(WHITE if name=="Verispect" else PURPLE),bold=True,align=PP_ALIGN.CENTER,space=4)
    txt(tf,price,18,col,bold=True,align=PP_ALIGN.CENTER,space=8)
    txt(tf,desc,13,col,align=PP_ALIGN.CENTER)
    x+=3.05

# 12 — GTM
s = slide(); title_only(s, "Go-to-market", "Founder-led, compliance-triggered, product-led-assisted")
t=box(s,0.8,2.0,11.7,4.3)
bullets(t.text_frame,[
    "Outbound to a tight ICP: AI-native startups in EU/SG/UAE using LLMs in high-risk decisions (start with HR-tech — probe library already fits).",
    "Free \"drift/bias snapshot\" lead magnet — one line, 5 min, returns a real report. Qualifies + delivers value.",
    "Compliance-led content + SEO — own \"EU AI Act monitoring\" and \"LLM bias detection.\"",
    "Land via free snapshot → expand Pro → graduate best logos to Enterprise (driven by their own audits/buyers).",
    "Channel: RAI consultancies, boutique law firms, GRC platforms as referral partners.",
], size=17, space=13)

# 13 — TRACTION / ROADMAP
s = slide(); title_only(s, "Plan — first 12 months")
t=box(s,0.8,2.0,11.7,4.3)
bullets(t.text_frame,[
    "Mo 0–3: ship multi-model + alerts + self-serve billing; 5 design partners; public launch.",
    "Mo 3–6: 15 paying Pro; first 1–2 Enterprise pilots; €3–5k MRR.",
    "Mo 6–12: 40+ Pro, 3–5 Enterprise; €15–25k MRR; SOC 2 Type I in progress; first hire.",
    "North-star metric: audit-ready reports generated per month.",
], size=18, space=16)

# 14 — TEAM
s = slide(); title_only(s, "Founder")
t=box(s,0.8,2.1,11.7,3.5)
txt(t.text_frame,"Ajay — Founder & Engineer",24,DARK,bold=True,space=10)
txt(t.text_frame,"BSCS, FAST NUCES Karachi. Built the entire stack solo: FastAPI proxy, privacy-preserving SDK, "
    "embedding drift engine, compliance PDF generator, React dashboard, auth.",17,GRAY,space=10)
txt(t.text_frame,"Found measurable hiring bias in a mainstream model in an afternoon — and decided someone should be watching continuously.",17,PURPLE,italic=True)

# 15 — CLOSE
s = slide(BG)
t=box(s,0.9,2.6,11.5,2.2)
txt(t.text_frame,"Your AI model changed.",36,WHITE,bold=True)
txt(t.text_frame,"We're the only ones who noticed.",36,PURPLE,bold=True)
u=box(s,0.9,5.0,11.5,1.4)
txt(u.text_frame,"Run a free snapshot — see your model's drift in 5 minutes.",18,RGBColor(0xC8,0xC3,0xF0))
txt(u.text_frame,"verispectai.com   ·   hello@verispectai.com",15,GRAY)

prs.save("Verispect-Pitch-Deck.pptx")
print("Saved Verispect-Pitch-Deck.pptx —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
